# Updated April 25: fixed slide 2 outline per Phase 1 feedback
"""
build_pptx.py
=============
CSCE 5222 – Group 18
Builds the final demo presentation by filling in the provided template.
"""

import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

TEMPLATE_PATH = '/home/ubuntu/upload/project2demoandslides.pptx'
OUTPUT_PATH   = '/home/ubuntu/project2/slides/Group18_Presentation.pptx'
FIG_DIR       = '/home/ubuntu/project2/figures'

prs = Presentation(TEMPLATE_PATH)

W = prs.slide_width
H = prs.slide_height

# ─────────────────────────────────────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────────────────────────────────────

def clear_slide_text(slide):
    """Remove all text from all text frames on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ''
            if shape.text_frame.paragraphs:
                shape.text_frame.paragraphs[0].runs[0].text = '' if shape.text_frame.paragraphs[0].runs else ''


def set_title(slide, title_text, font_size=28):
    """Set the title placeholder text."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title_text
            run.font.size = Pt(font_size)
            run.font.bold = True
            return


def set_body(slide, lines, font_size=18, placeholder_idx=1):
    """Set body text with bullet points."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            tf = shape.text_frame
            tf.clear()
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.level = 0
                if p.runs:
                    p.runs[0].font.size = Pt(font_size)
            return


def add_image(slide, img_path, left, top, width, height):
    """Add an image to a slide."""
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)


def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=None):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – Title Slide
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[0]
set_title(slide, 'Group 18: Feature Engineering for Time-Series', font_size=26)
# Update subtitle/body
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 1:
        tf = shape.text_frame
        tf.clear()
        lines = [
            'CSCE 5222 Feature Engineering',
            'Ashish Rathnakar Shetty (ID: 11808466)',
            'Kushal Sai Venigalla (ID: 11852559)',
            'Computer Science and Engineering Department',
            'University of North Texas',
        ]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            if p.runs:
                p.runs[0].font.size = Pt(16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – Outline
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[1]
set_title(slide, 'Outline', font_size=28)
set_body(slide, [
    'Part 1: Presentation (1–3 min)',
    '  • Recap of Project 1 Proposal',
    '  • Feature Engineering Techniques',
    '  • Machine Learning Models',
    '  • Results and Analysis',
    '  • Discussion & Conclusion',
    'Part 2: Live Demo (3–13 min)',
    '  • Dataset Exploration & Preprocessing',
    '  • Feature Extraction Demo',
    '  • Model Evaluation & Comparison',
], font_size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – Recap of Project 1 Proposal
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[2]
set_title(slide, 'Recap of Project 1 Proposal', font_size=26)
set_body(slide, [
    'Problem: Extracting meaningful features from time-series is a bottleneck.',
    'Motivation: Deep learning is powerful but slow and hard to interpret.',
    'We proposed a hybrid, automated, and explainable feature engineering framework.',
    'Five papers surveyed (2024–2025): COCALITE, SPLiT, PPG Anomaly, Lag Features, AutoFE-XAI.',
    'Three datasets targeted: UCR (classification), ETTh1 (forecasting), WESAD-like (anomaly).',
    'Key update since Phase 1: Implemented Catch22 from scratch in pure NumPy.',
], font_size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – Recap (cont.)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[3]
set_title(slide, 'Recap of Project 1 Proposal (cont.)', font_size=26)
set_body(slide, [
    'Phase 1 Feedback Addressed:',
    '  • Added Group Number (Group 18) to all submissions.',
    '  • Separated evaluations clearly by task (classification / forecasting / anomaly).',
    '  • Defined explicit baselines for each task.',
    '  • Completed the methodology flowchart (see Figure in report).',
    '  • Added detailed timeline with milestones.',
    '  • Discussed cross-domain generalization limits in the discussion section.',
], font_size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – Feature Engineering Techniques
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[4]
set_title(slide, 'Feature Engineering Techniques', font_size=26)
set_body(slide, [
    '1. Catch22 (22 canonical time-series features):',
    '   Autocorrelation, entropy, permutation entropy, Hurst exponent, spectral features...',
    '   Implemented from scratch in pure NumPy (no C-compiler dependency).',
    '2. Lag + Rolling Statistics:',
    '   Lag features (t-1 to t-24), rolling mean/std/min/max over windows of 3, 6, 12, 24.',
    '3. Self-Supervised MLP Embeddings:',
    '   MLP trained to predict x[t] from x[t-window:t]. Hidden layer = learned features.',
    '4. SHAP-Guided Feature Pruning:',
    '   LightGBM proxy → SHAP values → keep top 55-60% features by mean |SHAP|.',
], font_size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – Feature Engineering Techniques (cont.) – with flowchart
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[5]
set_title(slide, 'Feature Engineering Pipeline', font_size=26)
set_body(slide, ['Our complete pipeline from raw data to model evaluation:'], font_size=16)
add_image(slide,
          os.path.join(FIG_DIR, 'fig11_pipeline_flowchart.png'),
          Inches(0.3), Inches(1.5), Inches(9.0), Inches(5.0))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – Machine Learning Models
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[6]
set_title(slide, 'Machine Learning Models', font_size=26)
set_body(slide, [
    'Classification: LightGBM Classifier',
    '   Rationale: Fast gradient boosting, handles tabular features well, supports SHAP natively.',
    'Forecasting: LightGBM Regressor + Ridge Regression (baseline)',
    '   Rationale: LightGBM captures non-linear interactions in lag/rolling features.',
    'Anomaly Detection: Isolation Forest',
    '   Rationale: Unsupervised, no labeled data needed; works well with Catch22 features.',
    'Embedding Trainer: MLP Regressor (sklearn)',
    '   Rationale: Lightweight, interpretable hidden layer, fast to train.',
], font_size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – Machine Learning Models (cont.)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[7]
set_title(slide, 'Baseline vs. Proposed Models', font_size=26)
set_body(slide, [
    'Task 1 – Classification:',
    '   Baseline: LightGBM on raw time-series (first 50 time steps)',
    '   Proposed: LightGBM on Catch22 features → SHAP pruned',
    'Task 2 – Forecasting:',
    '   Baseline: Ridge Regression on lag-1 only',
    '   Proposed: LightGBM on Lag+Rolling+Embeddings → SHAP pruned',
    'Task 3 – Anomaly Detection:',
    '   Baseline: Isolation Forest on raw 50-sample windows',
    '   Proposed: Isolation Forest on Catch22 → SHAP-selected features',
], font_size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – Results and Analysis (Classification)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[8]
set_title(slide, 'Results: Classification (UCR Datasets)', font_size=26)
set_body(slide, ['Catch22 outperforms raw features on 3/5 datasets. SHAP pruning maintains accuracy.'], font_size=15)
add_image(slide,
          os.path.join(FIG_DIR, 'fig1_classification_accuracy.png'),
          Inches(0.2), Inches(1.5), Inches(9.2), Inches(4.8))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 – Discussion
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[9]
set_title(slide, 'Discussion', font_size=26)
set_body(slide, [
    'Key Insight 1 – SHAP as a decision tool, not just reporting:',
    '   Pruning 40-45% of features preserved accuracy across all tasks.',
    'Key Insight 2 – Catch22 generalizes well but not universally:',
    '   ItalyPowerDemand: raw features win (local shape patterns matter).',
    '   SyntheticControl: Catch22 wins (global statistics are discriminative).',
    'Key Insight 3 – Self-supervised embeddings underperform explicit statistics:',
    '   MLP embeddings (MSE=0.337) < Lag+Rolling (MSE=0.228) for ETTh1.',
    'Limitation: WESAD-like data is simulated; real wearable data may differ.',
], font_size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 – Conclusion
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[10]
set_title(slide, 'Conclusion', font_size=26)
set_body(slide, [
    'Feature engineering still matters, even in the deep learning era.',
    'Catch22 + LightGBM is a powerful, efficient alternative to end-to-end deep models.',
    'SHAP-guided pruning is a practical, principled way to reduce feature dimensionality.',
    'Lag and rolling statistics remain the strongest features for time-series forecasting.',
    'Future work: LLM-assisted feature suggestion, edge-device deployment, real WESAD data.',
], font_size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 – References
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[11]
set_title(slide, 'References', font_size=26)
set_body(slide, [
    '[1] H. Badi et al., "COCALITE," IEEE BigData 2024, pp. 5234–5243.',
    '[2] M. D\'Aversa et al., "SPLiT," IEEE BigData 2024, pp. 3412–3421.',
    '[3] L. Valerio et al., "Unsupervised Anomaly Detection in PPG," IEEE SENSORS 2024.',
    '[4] T. Okadome & J. Nakamura, "Lag Operation," IEEE Access, vol. 12, 2024.',
    '[5] O. Petrosian et al., "AutoFE-XAI," IEEE Access, vol. 13, 2025.',
    '[6] C. H. Lubba et al., "catch22," Data Mining & Knowledge Discovery, 2019.',
], font_size=14)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 – End
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[12]
set_title(slide, 'Thank You!', font_size=32)
set_body(slide, [
    'Group 18 – CSCE 5222 Feature Engineering',
    'Ashish Rathnakar Shetty | Kushal Sai Venigalla',
    '',
    'Questions welcome!',
], font_size=18)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 – Demo Setup
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[14]
set_title(slide, 'Part 2: Live Demo – Setup', font_size=26)
set_body(slide, [
    'Environment: Python 3.11, Jupyter Notebook / VS Code',
    'Key Libraries: NumPy, Pandas, LightGBM, SHAP, scikit-learn, Matplotlib',
    'Datasets pre-downloaded: UCR (GunPoint, ECG200), ETTh1.csv',
    'All dependencies installed and tested before the demo.',
    'Notebook: feature_engineering_pipeline.py (modular, annotated)',
], font_size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 – Flow of Demo
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[15]
set_title(slide, 'Flow of the Demo', font_size=26)
set_body(slide, [
    '1. Dataset Exploration: Show raw GunPoint time-series, ETTh1 columns.',
    '2. Feature Extraction: Run catch22_features() on a sample series live.',
    '3. SHAP Pruning: Show LightGBM + SHAP → feature importance bar chart.',
    '4. Model Evaluation: Compare accuracy/MSE before and after feature engineering.',
    '5. Forecasting: Plot actual vs. predicted oil temperature (ETTh1).',
    '6. Anomaly Detection: Show BVP signal with highlighted anomaly regions.',
], font_size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 – Interactivity
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[16]
set_title(slide, 'Interactivity', font_size=26)
set_body(slide, [
    'Live code walkthrough of the Catch22 feature extraction loop.',
    'Interactive SHAP bar chart: audience can ask about specific features.',
    'Real-time model comparison table printed to console.',
    'Visualization: forecast vs. actual plot generated live.',
    'Anomaly detection: adjust contamination parameter and observe F1 change.',
], font_size=16)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"Saved presentation to {OUTPUT_PATH}")
