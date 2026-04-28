"""
build_final_clean.py
====================
CSCE 5222 – Group 18
Single-pass builder: edits all 16 slides in one run, saves once.
Uses the UNT template for the logo and background.
Avoids all duplicate-name corruption by never adding then removing slides.

Strategy:
- Copy template to output
- Edit slides 1-13 in-place
- Slides 14-19 of the template are repurposed for our content (conclusion, refs, thanks)
  and the remaining 3 are cleared/hidden
- Save exactly once
"""

import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = '/home/ubuntu/upload/project2demoandslides.pptx'
OUTPUT   = '/home/ubuntu/project2/slides/Group18_UNT_Presentation.pptx'
FIG_DIR  = '/home/ubuntu/project2/figures'

shutil.copy2(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)

UNT_GREEN = RGBColor(0x00, 0x85, 0x3E)
BLACK     = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def clear_slide(slide):
    to_remove = [s for s in slide.shapes if not s.is_placeholder]
    for s in to_remove:
        s._element.getparent().remove(s._element)


def set_ph(slide, idx, lines):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, (text, fs, bold, color, level) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.level = level
                p.space_before = Pt(3)
                run = p.add_run()
                run.text = text
                if fs:   run.font.size  = Pt(fs)
                if bold is not None: run.font.bold = bold
                if color: run.font.color.rgb = color
            return ph


def add_tb(slide, text, left, top, width, height,
           fs=16, bold=False, color=BLACK,
           align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_img(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


def green_line(slide):
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(1.58), Inches(13.33), Inches(0.055))
    shape.fill.solid()
    shape.fill.fore_color.rgb = UNT_GREEN
    shape.line.fill.background()


def blank_slide(slide):
    """Clear all content from a slide (make it empty with just logo)."""
    clear_slide(slide)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx not in (12,):  # keep slide number
            try:
                ph.text_frame.clear()
                ph.text_frame.paragraphs[0].runs[0].text = '' if ph.text_frame.paragraphs[0].runs else ''
            except:
                pass


# ─── Get all 19 slides ────────────────────────────────────────────────────────
sl = list(prs.slides)  # 19 slides

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = sl[0]
clear_slide(s)
set_ph(s, 0, [
    ("Group 18: Feature Engineering for Time-Series", 32, False, BLACK, 0),
    ("Hybrid, Automated, and Interpretable Approaches", 24, False, DARK_GRAY, 0),
])
# Add names as explicit textboxes so they definitely show
add_tb(s, "CSCE 5222 Feature Engineering  |  Group 18",
       1.5, 3.9, 10.33, 0.5, 18, False, UNT_GREEN, PP_ALIGN.CENTER)
add_tb(s, "Ashish Rathnakar Shetty (ID: 11808466)  ·  Kushal Sai Venigalla (ID: 11852559)",
       1.5, 4.5, 10.33, 0.55, 20, True, BLACK, PP_ALIGN.CENTER)
add_tb(s, "Computer Science and Engineering Department  |  University of North Texas",
       1.5, 5.15, 10.33, 0.45, 16, False, DARK_GRAY, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OUTLINE
# ══════════════════════════════════════════════════════════════════════════════
s = sl[1]
clear_slide(s)
set_ph(s, 0, [("Outline", 36, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("Part 1 — Presentation (1–3 min)", 22, True, UNT_GREEN, 0),
    ("Recap of Project 1 Proposal", 19, False, BLACK, 1),
    ("Feature Engineering Techniques Implemented", 19, False, BLACK, 1),
    ("Machine Learning Models and Baselines", 19, False, BLACK, 1),
    ("Results and Analysis", 19, False, BLACK, 1),
    ("Discussion and Conclusion", 19, False, BLACK, 1),
    ("Part 2 — Live Demo (3–13 min)", 22, True, UNT_GREEN, 0),
    ("Dataset exploration and preprocessing walkthrough", 19, False, BLACK, 1),
    ("Live feature extraction and SHAP pruning demo", 19, False, BLACK, 1),
    ("Model evaluation: before vs. after feature engineering", 19, False, BLACK, 1),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RECAP OF PROJECT 1
# ══════════════════════════════════════════════════════════════════════════════
s = sl[2]
clear_slide(s)
set_ph(s, 0, [("Recap of Project 1 Proposal", 34, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("Problem: Extracting meaningful features from time-series is a bottleneck in ML pipelines.", 19, False, BLACK, 0),
    ("Motivation: Deep learning is powerful but computationally expensive and hard to interpret.", 19, False, BLACK, 0),
    ("We proposed a unified framework combining hybrid, automated, and explainable feature engineering.", 19, False, BLACK, 0),
    ("5 papers surveyed (2024–2025): COCALITE, SPLiT, PPG Anomaly, Lag Features, AutoFE-XAI.", 19, False, BLACK, 0),
    ("3 datasets: UCR Archive (classification), ETTh1 (forecasting), WESAD-like (anomaly detection).", 19, False, BLACK, 0),
    ("Phase 1 Feedback Addressed:", 20, True, UNT_GREEN, 0),
    ("Added Group Number 18 to all submissions", 18, False, BLACK, 1),
    ("Separated evaluations clearly by task", 18, False, BLACK, 1),
    ("Defined explicit baselines for each task", 18, False, BLACK, 1),
    ("Completed the methodology flowchart (was incomplete in Phase 1)", 18, False, BLACK, 1),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DATASETS
# ══════════════════════════════════════════════════════════════════════════════
s = sl[3]
clear_slide(s)
set_ph(s, 0, [("Datasets Used", 36, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("UCR Time Series Archive — Classification", 21, True, UNT_GREEN, 0),
    ("5 datasets: GunPoint, ECG200, ItalyPowerDemand, SyntheticControl, TwoLeadECG", 18, False, BLACK, 1),
    ("Covers motion, medical, energy, and synthetic domains", 18, False, BLACK, 1),
    ("ETTh1 (Electricity Transformer Temperature) — Forecasting", 21, True, UNT_GREEN, 0),
    ("17,420 hourly readings of 7 variables; target: Oil Temperature (OT)", 18, False, BLACK, 1),
    ("Real-world multivariate industrial benchmark dataset", 18, False, BLACK, 1),
    ("WESAD-like PPG Signal — Anomaly Detection", 21, True, UNT_GREEN, 0),
    ("Simulated blood volume pulse (BVP) signal with injected anomalies", 18, False, BLACK, 1),
    ("Anomaly types: amplitude spikes, flatlines, high-noise segments (~16% anomaly rate)", 18, False, BLACK, 1),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FEATURE ENGINEERING TECHNIQUES
# ══════════════════════════════════════════════════════════════════════════════
s = sl[4]
clear_slide(s)
set_ph(s, 0, [("Feature Engineering Techniques", 34, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("1.  Catch22 — 22 canonical time-series features", 21, True, UNT_GREEN, 0),
    ("Autocorrelation, entropy, Hurst exponent, spectral features — implemented in pure NumPy", 18, False, BLACK, 1),
    ("2.  Lag + Rolling Statistics", 21, True, UNT_GREEN, 0),
    ("Lag features (t-1 to t-24) + rolling mean, std, min, max over windows of 3, 6, 12, 24", 18, False, BLACK, 1),
    ("3.  Self-Supervised MLP Embeddings", 21, True, UNT_GREEN, 0),
    ("MLP trained to predict x[t] from x[t-24:t]; hidden layer activations = learned features", 18, False, BLACK, 1),
    ("4.  SHAP-Guided Feature Pruning", 21, True, UNT_GREEN, 0),
    ("LightGBM proxy → compute mean |SHAP| → keep top 55–60% of features", 18, False, BLACK, 1),
    ("XAI actively guides modeling decisions — not just post-hoc reporting", 18, False, BLACK, 1),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PIPELINE FLOWCHART
# ══════════════════════════════════════════════════════════════════════════════
s = sl[5]
clear_slide(s)
set_ph(s, 0, [("Methodology Pipeline", 36, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [("", 12, False, BLACK, 0)])
add_img(s, os.path.join(FIG_DIR, 'fig11_pipeline_flowchart.png'), 0.5, 1.7, 12.3, 5.6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODELS AND BASELINES
# ══════════════════════════════════════════════════════════════════════════════
s = sl[6]
clear_slide(s)
set_ph(s, 0, [("Machine Learning Models and Baselines", 30, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("Classification Task", 21, True, UNT_GREEN, 0),
    ("Baseline: LightGBM on raw time-series (first 50 time steps)", 18, False, BLACK, 1),
    ("Proposed: LightGBM on Catch22 → SHAP-pruned  |  Metrics: Accuracy, F1", 18, False, BLACK, 1),
    ("Forecasting Task", 21, True, UNT_GREEN, 0),
    ("Baseline: Ridge Regression on lag-1 feature only", 18, False, BLACK, 1),
    ("Proposed: LightGBM on Lag+Rolling+Embeddings → SHAP pruned  |  Metrics: MSE, MAE", 18, False, BLACK, 1),
    ("Anomaly Detection Task", 21, True, UNT_GREEN, 0),
    ("Baseline: Isolation Forest on raw 50-sample windows", 18, False, BLACK, 1),
    ("Proposed: Isolation Forest on Catch22 → SHAP-selected  |  Metrics: F1, AUC", 18, False, BLACK, 1),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS: CLASSIFICATION
# ══════════════════════════════════════════════════════════════════════════════
s = sl[7]
clear_slide(s)
set_ph(s, 0, [("Results: Classification (UCR Datasets)", 30, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [("", 12, False, BLACK, 0)])
add_img(s, os.path.join(FIG_DIR, 'fig1_classification_accuracy.png'), 0.4, 1.7, 8.5, 4.5)
add_tb(s, "Key Findings:", 9.1, 1.8, 4.0, 0.4, 17, True, UNT_GREEN)
for i, f in enumerate([
    "SyntheticControl: Catch22 +6.3%",
    "GunPoint: Catch22 +6.0%",
    "ItalyPowerDemand: Raw wins",
    "(local shape patterns matter)",
    "SHAP pruning: 9/22 features",
    "removed, zero accuracy loss",
    "TwoLeadECG: tied (23 train samples)",
]):
    add_tb(s, f"• {f}", 9.1, 2.3 + i * 0.5, 4.0, 0.45, 15, False, BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SHAP IMPORTANCE
# ══════════════════════════════════════════════════════════════════════════════
s = sl[8]
clear_slide(s)
set_ph(s, 0, [("SHAP Feature Importance — GunPoint Dataset", 28, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [("", 12, False, BLACK, 0)])
add_img(s, os.path.join(FIG_DIR, 'fig3_shap_importance.png'), 0.4, 1.7, 7.5, 5.5)
add_tb(s, "How SHAP Guides Decisions:", 8.1, 1.8, 5.0, 0.4, 17, True, UNT_GREEN)
for i, (t, c) in enumerate([
    ("Top: acf_lag1, diff_variance,", BLACK),
    ("hurst, acf_lag2", BLACK),
    ("→ temporal dependence & roughness", UNT_GREEN),
    ("", BLACK),
    ("Pruned: hist_entropy, peak_freq", BLACK),
    ("→ removed without accuracy loss", UNT_GREEN),
    ("", BLACK),
    ("22 → 13 features (40% reduction)", BLACK),
    ("XAI as a design tool,", BLACK),
    ("not just post-hoc reporting", BLACK),
]):
    add_tb(s, t, 8.1, 2.3 + i * 0.44, 5.0, 0.4, 15, False, c)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RESULTS: FORECASTING
# ══════════════════════════════════════════════════════════════════════════════
s = sl[9]
clear_slide(s)
set_ph(s, 0, [("Results: Forecasting (ETTh1 — Oil Temperature)", 28, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [("", 12, False, BLACK, 0)])
add_img(s, os.path.join(FIG_DIR, 'fig4_forecasting_metrics.png'), 0.4, 1.7, 7.8, 3.0)
add_img(s, os.path.join(FIG_DIR, 'fig5_forecast_vs_actual.png'), 0.4, 4.8, 7.8, 2.5)
add_tb(s, "Key Results:", 8.4, 1.8, 4.7, 0.4, 17, True, UNT_GREEN)
for i, (t, bold, c) in enumerate([
    ("Raw (lag-1):     MSE = 1.285", False, BLACK),
    ("Lag+Rolling:     MSE = 0.228  (−82%)", True, UNT_GREEN),
    ("+ Embeddings:    MSE = 0.337", False, BLACK),
    ("+ SHAP Prune:    MSE = 0.230  (−82%)", True, UNT_GREEN),
    ("", False, BLACK),
    ("SHAP removed 12/28 features", False, BLACK),
    ("with no MSE increase", False, BLACK),
    ("Self-supervised embeddings did", False, BLACK),
    ("not outperform explicit statistics", False, BLACK),
]):
    add_tb(s, t, 8.4, 2.3 + i * 0.5, 4.7, 0.45, 15, bold, c)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS: ANOMALY DETECTION
# ══════════════════════════════════════════════════════════════════════════════
s = sl[10]
clear_slide(s)
set_ph(s, 0, [("Results: Anomaly Detection (WESAD-like PPG)", 28, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [("", 12, False, BLACK, 0)])
add_img(s, os.path.join(FIG_DIR, 'fig6_anomaly_signal.png'), 0.4, 1.7, 8.5, 2.5)
add_img(s, os.path.join(FIG_DIR, 'fig7_anomaly_metrics.png'), 0.4, 4.3, 8.5, 2.9)
add_tb(s, "Key Results:", 9.1, 1.8, 4.0, 0.4, 17, True, UNT_GREEN)
for i, (t, bold, c) in enumerate([
    ("Raw Windows:   F1 = 0.871", False, BLACK),
    ("Catch22 Only:  F1 = 0.581 ↓", False, DARK_GRAY),
    ("Catch22+SHAP:  F1 = 0.710 ↑", True, UNT_GREEN),
    ("", False, BLACK),
    ("Why Catch22 hurt:", True, BLACK),
    ("Amplitude normalization", False, BLACK),
    ("masks spike anomalies", False, BLACK),
    ("SHAP selected diff_variance", False, BLACK),
    ("& time_reversibility to", False, BLACK),
    ("recover F1 score", False, BLACK),
]):
    add_tb(s, t, 9.1, 2.3 + i * 0.5, 4.0, 0.45, 15, bold, c)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SUMMARY VISUALIZATIONS
# ══════════════════════════════════════════════════════════════════════════════
s = sl[11]
clear_slide(s)
set_ph(s, 0, [("Results Summary and Visualizations", 32, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [("", 12, False, BLACK, 0)])
add_img(s, os.path.join(FIG_DIR, 'fig12_summary_heatmap.png'), 0.4, 1.7, 7.5, 3.6)
add_img(s, os.path.join(FIG_DIR, 'fig9_pca_projection.png'), 8.1, 1.7, 5.0, 3.6)
add_img(s, os.path.join(FIG_DIR, 'fig8_feature_correlation.png'), 0.4, 5.4, 5.5, 1.9)
add_img(s, os.path.join(FIG_DIR, 'fig10_forecasting_shap.png'), 6.1, 5.4, 7.0, 1.9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DISCUSSION
# ══════════════════════════════════════════════════════════════════════════════
s = sl[12]
clear_slide(s)
set_ph(s, 0, [("Discussion", 36, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("When Catch22 Works Best", 21, True, UNT_GREEN, 0),
    ("Datasets where global statistics are discriminative (SyntheticControl, GunPoint)", 18, False, BLACK, 1),
    ("Compresses 150+ time steps into 22 interpretable features efficiently", 18, False, BLACK, 1),
    ("When Raw Features Win", 21, True, UNT_GREEN, 0),
    ("ItalyPowerDemand: class defined by a sharp drop at a specific timestamp", 18, False, BLACK, 1),
    ("Catch22 averages out local patterns — global features miss local timing", 18, False, BLACK, 1),
    ("SHAP as an Active Design Tool", 21, True, UNT_GREEN, 0),
    ("Reduced feature space by 40–45% across all three tasks", 18, False, BLACK, 1),
    ("Faster inference, lower memory, more interpretable — critical for edge deployment", 18, False, BLACK, 1),
    ("Challenges: WESAD-like data is simulated; real wearable data may differ significantly", 18, False, DARK_GRAY, 0),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSION (repurpose template slide 14)
# ══════════════════════════════════════════════════════════════════════════════
s = sl[13]
clear_slide(s)
set_ph(s, 0, [("Conclusion and Future Work", 34, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("Key Conclusions", 21, True, UNT_GREEN, 0),
    ("Catch22 + LightGBM matches deep learning accuracy at a fraction of the cost", 18, False, BLACK, 1),
    ("Lag + Rolling Statistics are the strongest forecasting features (MSE −82%)", 18, False, BLACK, 1),
    ("SHAP-guided pruning reliably reduces dimensionality with zero accuracy loss", 18, False, BLACK, 1),
    ("Cross-domain generalization is limited when local timing is the discriminant", 18, False, BLACK, 1),
    ("Future Work", 21, True, UNT_GREEN, 0),
    ("LLM-assisted feature suggestion: use LLMs to propose domain-specific features", 18, False, BLACK, 1),
    ("Edge device deployment: ultra-lightweight Catch22 for real-time IoT sensors", 18, False, BLACK, 1),
    ("Validate on real WESAD wearable stress dataset (requires institutional access)", 18, False, BLACK, 1),
    ("XAI-first AutoML: make SHAP a primary optimization objective", 18, False, BLACK, 1),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — REFERENCES (repurpose template slide 15)
# ══════════════════════════════════════════════════════════════════════════════
s = sl[14]
clear_slide(s)
set_ph(s, 0, [("References", 36, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("[1] H. Badi, G. Forestier, and J. Weber, \"COCALITE: A Hybrid Approach for Time Series Classification Using Catch22 and LITE,\" in 2024 IEEE BigData, pp. 5234–5243.", 15, False, BLACK, 0),
    ("[2] M. D'Aversa, F. Giannini, and M. Scardapane, \"SPLiT: Spatio-Temporal Linear Model Trees for Multi-Step Time Series Forecasting,\" in 2024 IEEE BigData, pp. 3412–3421.", 15, False, BLACK, 0),
    ("[3] L. Valerio, A. Passarella, and M. Conti, \"Unsupervised Anomaly Detection in Wearable PPG Data via Catch22 Features,\" in 2024 IEEE SENSORS, pp. 1–4.", 15, False, BLACK, 0),
    ("[4] T. Okadome and J. Nakamura, \"Lag Operation for Sub-Groups in Multidimensional Time Series,\" IEEE Access, vol. 12, pp. 15432–15445, 2024.", 15, False, BLACK, 0),
    ("[5] O. Petrosian, D. Wang, and S. Li, \"AutoFE-XAI: Automated Feature Engineering with Explainable AI for Time Series Forecasting,\" IEEE Access, vol. 13, pp. 2109–2120, 2025.", 15, False, BLACK, 0),
    ("[6] C. H. Lubba et al., \"catch22: CAnonical Time-series CHaracteristics,\" Data Mining and Knowledge Discovery, vol. 33, no. 6, pp. 1821–1852, 2019.", 15, False, BLACK, 0),
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU (repurpose template slide 16)
# ══════════════════════════════════════════════════════════════════════════════
s = sl[15]
clear_slide(s)
set_ph(s, 0, [("Thank You!", 48, False, UNT_GREEN, 0)])
add_tb(s, "Questions & Discussion",
       1.5, 3.0, 10.33, 0.55, 24, False, DARK_GRAY, PP_ALIGN.CENTER)
add_tb(s, "Ashish Rathnakar Shetty  ·  ashishrathnakarshetty@unt.edu",
       1.5, 3.7, 10.33, 0.45, 17, False, BLACK, PP_ALIGN.CENTER)
add_tb(s, "Kushal Sai Venigalla  ·  KushalSaiVenigalla@unt.edu",
       1.5, 4.2, 10.33, 0.45, 17, False, BLACK, PP_ALIGN.CENTER)
add_tb(s, "GitHub: github.com/Ash-projects-personal/CSCE5222-Group18-FeatureEngineering",
       1.5, 4.85, 10.33, 0.45, 15, False, UNT_GREEN, PP_ALIGN.CENTER)
add_tb(s, "CSCE 5222 Feature Engineering  |  Group 18  |  University of North Texas",
       1.5, 5.4, 10.33, 0.4, 14, False, DARK_GRAY, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDES 17-19 — Clear leftover template slides
# ══════════════════════════════════════════════════════════════════════════════
for s in sl[16:]:
    blank_slide(s)

# ─── Remove slides 17-19 from the presentation ────────────────────────────────
xml_slides = prs.slides._sldIdLst
all_ids = list(xml_slides)
for sld_id in all_ids[16:]:
    xml_slides.remove(sld_id)

print(f"Final slide count: {len(prs.slides)}")

# ─── Save once ────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
