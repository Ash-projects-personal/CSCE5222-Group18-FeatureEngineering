"""
build_final_pptx.py
===================
CSCE 5222 – Group 18
Builds a completely original, polished presentation from scratch.
No template slides included — pure content only.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FIG_DIR = '/home/ubuntu/project2/figures'
OUTPUT  = '/home/ubuntu/project2/slides/Group18_Final_Presentation.pptx'

# ─── Color palette ────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x37, 0x5E)   # deep navy
MID_BLUE    = RGBColor(0x2E, 0x6D, 0xA4)   # medium blue
ACCENT      = RGBColor(0xE8, 0x6A, 0x2B)   # orange accent
LIGHT_GRAY  = RGBColor(0xF4, 0xF6, 0xF9)   # slide bg
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x1C, 0x1C, 0x1C)
TEXT_MID    = RGBColor(0x44, 0x44, 0x44)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
RED_SOFT    = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank layout

# ─── Helper functions ─────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=16, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, space_before=6, indent=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_image(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width), Inches(height))


def slide_header(slide, title, subtitle=None):
    """Add a consistent top header bar to every content slide."""
    add_rect(slide, 0, 0, 13.33, 1.1, DARK_BLUE)
    add_text(slide, title, 0.35, 0.12, 10, 0.75,
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.72, 10, 0.35,
                 font_size=13, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8),
                 align=PP_ALIGN.LEFT)
    # Accent bar
    add_rect(slide, 0, 1.1, 13.33, 0.06, ACCENT)
    # Slide background
    add_rect(slide, 0, 1.16, 13.33, 6.34, LIGHT_GRAY)
    # Footer
    add_rect(slide, 0, 7.15, 13.33, 0.35, DARK_BLUE)
    add_text(slide, "CSCE 5222 Feature Engineering  |  Group 18  |  Ashish Rathnakar Shetty & Kushal Sai Venigalla",
             0.2, 7.17, 12, 0.28, font_size=9, color=RGBColor(0xB0, 0xC8, 0xE8),
             align=PP_ALIGN.LEFT)


def content_box(slide, left, top, width, height, title=None, bg=WHITE):
    """White card with optional title."""
    shape = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(0.5)
    if title:
        add_text(slide, title, left + 0.12, top + 0.08, width - 0.2, 0.38,
                 font_size=13, bold=True, color=MID_BLUE)
    return shape


def bullet_box(slide, items, left, top, width, height,
               title=None, font_size=15, bullet="▸"):
    content_box(slide, left, top, width, height, title)
    y_start = top + (0.55 if title else 0.15)
    for item in items:
        if item.startswith("  "):
            txt = f"  ◦ {item.strip()}"
            fs = font_size - 1.5
            col = TEXT_MID
        else:
            txt = f"{bullet}  {item}"
            fs = font_size
            col = TEXT_DARK
        add_text(slide, txt, left + 0.15, y_start, width - 0.3, 0.38,
                 font_size=fs, color=col, wrap=True)
        y_start += 0.38


def metric_card(slide, left, top, width, height, label, value, delta=None,
                val_color=None):
    """A KPI card with big number."""
    content_box(slide, left, top, width, height)
    add_text(slide, label, left + 0.1, top + 0.1, width - 0.2, 0.35,
             font_size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)
    vc = val_color or MID_BLUE
    add_text(slide, value, left + 0.05, top + 0.42, width - 0.1, 0.55,
             font_size=26, bold=True, color=vc, align=PP_ALIGN.CENTER)
    if delta:
        add_text(slide, delta, left + 0.05, top + 0.95, width - 0.1, 0.3,
                 font_size=11, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# Full background gradient effect
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0x0D, 0x1F, 0x3C))

# Left accent strip
add_rect(slide, 0, 0, 0.55, 7.5, ACCENT)

# Title
add_text(slide,
    "Feature Engineering for Time-Series",
    0.85, 1.4, 11.5, 1.2,
    font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_text(slide,
    "Hybrid, Automated, and Interpretable Approaches",
    0.85, 2.65, 11.0, 0.6,
    font_size=22, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8),
    align=PP_ALIGN.LEFT)

# Divider line
add_rect(slide, 0.85, 3.35, 7.5, 0.05, ACCENT)

# Authors
add_text(slide,
    "Ashish Rathnakar Shetty  ·  Kushal Sai Venigalla",
    0.85, 3.55, 10, 0.5,
    font_size=17, color=WHITE, align=PP_ALIGN.LEFT)

add_text(slide,
    "Department of Computer Science and Engineering\nUniversity of North Texas",
    0.85, 4.1, 10, 0.65,
    font_size=14, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)

add_text(slide,
    "CSCE 5222 Feature Engineering  |  Group 18",
    0.85, 4.85, 10, 0.45,
    font_size=13, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.LEFT)

# Bottom bar
add_rect(slide, 0, 6.9, 13.33, 0.6, RGBColor(0x0D, 0x1F, 0x3C))
add_text(slide, "Project 2 — Implementation & Full Report",
         0.85, 6.95, 12, 0.45,
         font_size=12, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OUTLINE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Outline", "What we'll cover in this presentation")

sections = [
    ("01", "Recap of Project 1 Proposal",        "Problem, motivation, and what we planned"),
    ("02", "Feature Engineering Techniques",      "Catch22, Lag+Rolling, Self-Supervised, SHAP"),
    ("03", "Machine Learning Models & Baselines", "LightGBM, Ridge, Isolation Forest"),
    ("04", "Results and Analysis",                "Classification · Forecasting · Anomaly Detection"),
    ("05", "Discussion",                          "Key insights, challenges, and limitations"),
    ("06", "Conclusion & Future Work",            "Takeaways and next steps"),
]

for i, (num, title, desc) in enumerate(sections):
    col = i % 3
    row = i // 3
    lx = 0.35 + col * 4.3
    ty = 1.4 + row * 2.8

    add_rect(slide, lx, ty, 4.0, 2.4, WHITE)
    # Number badge
    add_rect(slide, lx, ty, 0.7, 0.55, MID_BLUE)
    add_text(slide, num, lx + 0.02, ty + 0.04, 0.66, 0.45,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, lx + 0.08, ty + 0.65, 3.8, 0.55,
             font_size=14, bold=True, color=DARK_BLUE)
    add_text(slide, desc, lx + 0.08, ty + 1.25, 3.8, 0.9,
             font_size=12, color=TEXT_MID, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RECAP OF PROJECT 1
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Recap: Project 1 Proposal",
             "What we proposed and how Phase 2 builds on it")

# Left column — problem & motivation
bullet_box(slide, [
    "Time-series data needs specialized feature engineering",
    "Deep learning is powerful but slow and uninterpretable",
    "AutoFE tools generate thousands of redundant features",
    "No unified framework balancing speed, accuracy, and explainability",
], 0.35, 1.3, 5.9, 2.8, title="Problem & Motivation", font_size=14)

# Right column — what we proposed
bullet_box(slide, [
    "Survey 5 recent papers (2024–2025)",
    "Propose hybrid Catch22 + LightGBM pipeline",
    "Automate feature selection via SHAP values",
    "Test on UCR, ETTh1, and WESAD-like datasets",
], 6.5, 1.3, 6.5, 2.8, title="Phase 1 Proposal", font_size=14)

# Phase 1 feedback addressed
add_rect(slide, 0.35, 4.3, 12.6, 2.55, RGBColor(0xE8, 0xF4, 0xFF))
add_text(slide, "✓  Phase 1 Feedback Addressed in Phase 2",
         0.55, 4.38, 12, 0.4, font_size=13, bold=True, color=MID_BLUE)

feedback = [
    "Added Group Number 18 to all submissions",
    "Separated evaluations by task (classification / forecasting / anomaly detection)",
    "Defined explicit baselines for each task",
    "Completed the methodology flowchart (was incomplete in Phase 1)",
    "Added detailed timeline with milestones",
    "Discussed cross-domain generalization limits explicitly",
]
col1 = feedback[:3]
col2 = feedback[3:]
for i, item in enumerate(col1):
    add_text(slide, f"▸  {item}", 0.55, 4.85 + i*0.38, 6.0, 0.35,
             font_size=12, color=GREEN)
for i, item in enumerate(col2):
    add_text(slide, f"▸  {item}", 6.7, 4.85 + i*0.38, 6.0, 0.35,
             font_size=12, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DATASETS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Datasets Used",
             "Three diverse datasets covering classification, forecasting, and anomaly detection")

datasets = [
    ("UCR Time Series Archive",
     "Classification",
     ["5 datasets: GunPoint, ECG200, ItalyPowerDemand,", "SyntheticControl, TwoLeadECG",
      "Covers motion, medical, energy, and synthetic domains",
      "Train sizes: 23–300 samples per dataset"],
     "📊"),
    ("ETTh1 — Electricity Transformer",
     "Forecasting",
     ["17,420 hourly readings of 7 sensor variables",
      "Target: Oil Temperature (OT)",
      "Real-world multivariate industrial data",
      "Standard benchmark for time-series forecasting"],
     "⚡"),
    ("WESAD-like PPG Signal",
     "Anomaly Detection",
     ["Simulated blood volume pulse (BVP) signal",
      "Injected anomalies: spikes, flatlines, high noise",
      "2,000 samples, ~16% anomaly rate",
      "Models wearable sensor artifact detection"],
     "❤️"),
]

for i, (name, task, bullets, icon) in enumerate(datasets):
    lx = 0.35 + i * 4.32
    add_rect(slide, lx, 1.35, 4.1, 5.55, WHITE)
    # Task badge
    add_rect(slide, lx, 1.35, 4.1, 0.45, MID_BLUE)
    add_text(slide, task, lx + 0.1, 1.38, 3.9, 0.38,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, name, lx + 0.12, 1.9, 3.85, 0.55,
             font_size=14, bold=True, color=DARK_BLUE)
    for j, b in enumerate(bullets):
        add_text(slide, f"▸  {b}", lx + 0.12, 2.55 + j * 0.52, 3.85, 0.48,
                 font_size=12, color=TEXT_MID, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FEATURE ENGINEERING TECHNIQUES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Feature Engineering Techniques",
             "Four complementary strategies implemented and evaluated")

techniques = [
    ("Catch22", MID_BLUE,
     ["22 canonical time-series features",
      "Autocorrelation, entropy, Hurst exponent,",
      "spectral features, permutation entropy",
      "Implemented from scratch in pure NumPy"]),
    ("Lag + Rolling Stats", RGBColor(0x16, 0x7A, 0x5E),
     ["Lag features: t-1, t-2, ..., t-24",
      "Rolling mean, std, min, max",
      "Windows: 3, 6, 12, 24 time steps",
      "Captures local temporal dynamics"]),
    ("Self-Supervised Embeddings", RGBColor(0x7D, 0x3C, 0x98),
     ["MLP trained to predict x[t] from x[t-24:t]",
      "Hidden layer activations = learned features",
      "No labels needed — fully unsupervised",
      "Captures non-linear temporal patterns"]),
    ("SHAP-Guided Pruning", ACCENT,
     ["LightGBM proxy → compute SHAP values",
      "Keep top 55–60% features by mean |SHAP|",
      "Removes redundant/noisy features",
      "XAI actively guides modeling decisions"]),
]

for i, (name, color, bullets) in enumerate(techniques):
    lx = 0.35 + i * 3.24
    add_rect(slide, lx, 1.35, 3.1, 5.55, WHITE)
    add_rect(slide, lx, 1.35, 3.1, 0.5, color)
    add_text(slide, name, lx + 0.1, 1.38, 2.9, 0.42,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bullets):
        add_text(slide, f"▸  {b}", lx + 0.1, 2.0 + j * 0.52, 2.9, 0.48,
                 font_size=11.5, color=TEXT_MID, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PIPELINE FLOWCHART
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Methodology Pipeline",
             "End-to-end flow from raw data to evaluation")

add_image(slide, os.path.join(FIG_DIR, 'fig11_pipeline_flowchart.png'),
          0.5, 1.3, 12.3, 5.9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODELS & BASELINES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Machine Learning Models & Baselines",
             "Clear baselines defined for each task to quantify feature engineering gains")

rows = [
    ("Classification", "LightGBM on raw time-series\n(first 50 time steps)",
     "LightGBM on Catch22 features\n→ SHAP-pruned subset",
     "Accuracy, Weighted F1-Score"),
    ("Forecasting", "Ridge Regression on\nlag-1 feature only",
     "LightGBM on Lag+Rolling+\nSelf-Supervised Embeddings",
     "MSE, MAE"),
    ("Anomaly Detection", "Isolation Forest on\nraw 50-sample windows",
     "Isolation Forest on\nCatch22 → SHAP-selected features",
     "F1-Score, ROC-AUC"),
]

# Table header
add_rect(slide, 0.35, 1.35, 12.6, 0.5, DARK_BLUE)
for j, hdr in enumerate(["Task", "Baseline", "Proposed Method", "Metrics"]):
    add_text(slide, hdr, 0.45 + j * 3.15, 1.4, 3.0, 0.4,
             font_size=13, bold=True, color=WHITE)

for i, (task, baseline, proposed, metrics) in enumerate(rows):
    bg = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF5, 0xFB)
    add_rect(slide, 0.35, 1.85 + i * 1.6, 12.6, 1.55, bg)
    add_text(slide, task, 0.45, 1.9 + i * 1.6, 3.0, 1.4,
             font_size=14, bold=True, color=DARK_BLUE, wrap=True)
    add_text(slide, baseline, 3.6, 1.9 + i * 1.6, 3.0, 1.4,
             font_size=12, color=RED_SOFT, wrap=True)
    add_text(slide, proposed, 6.75, 1.9 + i * 1.6, 3.0, 1.4,
             font_size=12, color=GREEN, wrap=True)
    add_text(slide, metrics, 9.9, 1.9 + i * 1.6, 3.0, 1.4,
             font_size=12, color=TEXT_MID, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS: CLASSIFICATION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Results: Time-Series Classification",
             "UCR Archive — 5 datasets — LightGBM with Raw vs. Catch22 vs. Catch22+SHAP")

# Chart on left
add_image(slide, os.path.join(FIG_DIR, 'fig1_classification_accuracy.png'),
          0.35, 1.3, 7.8, 4.5)

# Key findings on right
add_rect(slide, 8.4, 1.3, 4.6, 4.5, WHITE)
add_text(slide, "Key Findings", 8.55, 1.38, 4.3, 0.42,
         font_size=14, bold=True, color=MID_BLUE)

findings = [
    ("SyntheticControl", "+6.3%", "Catch22 wins — global stats discriminative"),
    ("GunPoint", "+6.0%", "Catch22 wins — motion pattern captured"),
    ("ItalyPowerDemand", "−21%", "Raw wins — local shape patterns matter"),
    ("SHAP Pruning", "0% loss", "Removed 9/22 features, accuracy unchanged"),
    ("TwoLeadECG", "Tied", "Too few training samples (23 only)"),
]

for i, (ds, delta, note) in enumerate(findings):
    y = 1.9 + i * 0.75
    col = GREEN if "+" in delta or "0%" in delta else (RED_SOFT if "−" in delta else TEXT_MID)
    add_text(slide, ds, 8.55, y, 1.8, 0.35, font_size=12, bold=True, color=DARK_BLUE)
    add_text(slide, delta, 10.35, y, 0.9, 0.35, font_size=13, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(slide, note, 8.55, y + 0.32, 4.3, 0.35, font_size=10.5, color=TEXT_MID, wrap=True)

# Bottom summary bar
add_rect(slide, 0.35, 5.95, 12.6, 0.85, RGBColor(0xE8, 0xF4, 0xFF))
add_text(slide,
    "▸  Catch22 + LightGBM matches or beats raw features on 3/5 datasets with only 22 compact features",
    0.55, 6.05, 12.2, 0.65, font_size=13, color=DARK_BLUE, bold=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — RESULTS: SHAP FEATURE IMPORTANCE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "SHAP Feature Importance",
             "SHAP values actively guide feature selection — not just post-hoc explanation")

add_image(slide, os.path.join(FIG_DIR, 'fig3_shap_importance.png'),
          0.35, 1.3, 6.5, 5.9)

# Right side explanation
add_rect(slide, 7.1, 1.3, 5.9, 5.9, WHITE)
add_text(slide, "How SHAP Guides Modeling", 7.25, 1.38, 5.6, 0.42,
         font_size=14, bold=True, color=MID_BLUE)

shap_points = [
    "Top features (GunPoint): acf_lag1, diff_variance, hurst, acf_lag2",
    "These capture temporal dependence and roughness — exactly what distinguishes gun-draw from pointing",
    "Bottom features: hist_entropy, peak_freq — pruned away without accuracy loss",
    "SHAP reduced 22 features → 13 features (40% reduction)",
    "Same accuracy, faster inference, more interpretable model",
    "This is XAI being used as a design tool, not just a reporting tool",
]
for i, pt in enumerate(shap_points):
    add_text(slide, f"▸  {pt}", 7.25, 1.95 + i * 0.82, 5.6, 0.75,
             font_size=12, color=TEXT_DARK, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RESULTS: FORECASTING
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Results: Multivariate Forecasting (ETTh1)",
             "Predicting Oil Temperature — 4 methods compared")

# Metric cards at top
cards = [
    ("Raw (Lag-1)", "MSE 1.285", None, RED_SOFT),
    ("Lag + Rolling", "MSE 0.228", "−82%", GREEN),
    ("+ Embeddings", "MSE 0.337", "−74%", MID_BLUE),
    ("+ SHAP Pruning", "MSE 0.230", "−82%", GREEN),
]
for i, (label, val, delta, col) in enumerate(cards):
    metric_card(slide, 0.35 + i * 3.24, 1.35, 3.1, 1.55,
                label, val, delta, col)

# Chart
add_image(slide, os.path.join(FIG_DIR, 'fig5_forecast_vs_actual.png'),
          0.35, 3.1, 8.5, 3.65)

# Right findings
add_rect(slide, 9.1, 3.1, 4.0, 3.65, WHITE)
add_text(slide, "Key Insights", 9.25, 3.18, 3.7, 0.4,
         font_size=13, bold=True, color=MID_BLUE)
insights = [
    "Lag+Rolling reduces MSE by 82% over raw baseline",
    "Self-supervised embeddings add no benefit over explicit stats",
    "SHAP removes 12/28 features with zero MSE increase",
    "Lag-1 and Lag-2 are the most important features (SHAP)",
    "Rolling mean over 24h is the top window feature",
]
for i, ins in enumerate(insights):
    add_text(slide, f"▸  {ins}", 9.25, 3.65 + i * 0.6, 3.7, 0.55,
             font_size=11.5, color=TEXT_DARK, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS: ANOMALY DETECTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Results: Anomaly Detection (WESAD-like PPG)",
             "Unsupervised detection of signal artifacts using Isolation Forest")

# Signal plot
add_image(slide, os.path.join(FIG_DIR, 'fig6_anomaly_signal.png'),
          0.35, 1.3, 8.2, 2.9)

# Metrics chart
add_image(slide, os.path.join(FIG_DIR, 'fig7_anomaly_metrics.png'),
          0.35, 4.3, 5.5, 2.9)

# Right panel
add_rect(slide, 8.75, 1.3, 4.25, 5.9, WHITE)
add_text(slide, "Key Findings", 8.9, 1.38, 4.0, 0.4,
         font_size=14, bold=True, color=MID_BLUE)

anom_findings = [
    ("Raw Windows", "F1=0.871", GREEN, "Spikes alter amplitude directly — easy to detect"),
    ("Catch22 Only", "F1=0.581", RED_SOFT, "Amplitude normalization masks spike anomalies"),
    ("Catch22+SHAP", "F1=0.710", MID_BLUE, "SHAP selects variance/roughness features that survive normalization"),
]
y = 1.9
for method, score, col, note in anom_findings:
    add_text(slide, method, 8.9, y, 2.1, 0.35, font_size=12, bold=True, color=DARK_BLUE)
    add_text(slide, score, 11.0, y, 1.8, 0.35, font_size=13, bold=True, color=col)
    add_text(slide, note, 8.9, y + 0.35, 4.0, 0.45, font_size=10.5, color=TEXT_MID, wrap=True)
    y += 0.95

add_rect(slide, 8.75, 5.0, 4.25, 2.2, RGBColor(0xF0, 0xF5, 0xFB))
add_text(slide, "Why did Catch22 hurt?", 8.9, 5.08, 4.0, 0.38,
         font_size=12, bold=True, color=DARK_BLUE)
add_text(slide,
    "Catch22 normalizes the signal, removing absolute amplitude information. "
    "Spike anomalies are defined by high amplitude — so normalization erases the signal. "
    "SHAP identified diff_variance and time_reversibility as the features that still "
    "capture flatline and noise anomalies after normalization.",
    8.9, 5.5, 4.0, 1.6, font_size=11, color=TEXT_MID, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SUMMARY TABLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Results Summary",
             "All tasks and methods at a glance")

add_image(slide, os.path.join(FIG_DIR, 'fig12_summary_heatmap.png'),
          0.35, 1.3, 7.5, 5.9)

# Right: PCA projection
add_rect(slide, 8.1, 1.3, 5.0, 2.8, WHITE)
add_text(slide, "PCA of Catch22 Features", 8.25, 1.38, 4.7, 0.4,
         font_size=13, bold=True, color=MID_BLUE)
add_image(slide, os.path.join(FIG_DIR, 'fig9_pca_projection.png'),
          8.1, 1.8, 5.0, 2.3)

# Right: correlation heatmap
add_rect(slide, 8.1, 4.3, 5.0, 2.8, WHITE)
add_text(slide, "Feature Correlation Heatmap", 8.25, 4.38, 4.7, 0.4,
         font_size=13, bold=True, color=MID_BLUE)
add_image(slide, os.path.join(FIG_DIR, 'fig8_feature_correlation.png'),
          8.1, 4.8, 5.0, 2.3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DISCUSSION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Discussion",
             "Interpreting results, challenges faced, and limitations")

# Three discussion panels
panels = [
    ("When Catch22 Wins", MID_BLUE,
     ["Datasets where global statistics are discriminative (SyntheticControl, GunPoint)",
      "Catch22 compresses 150+ time steps into 22 interpretable numbers",
      "Works well when class differences are in distribution shape, not local timing"]),
    ("When Raw Features Win", RED_SOFT,
     ["ItalyPowerDemand: class defined by a sharp drop at a specific timestamp",
      "Catch22 averages out local patterns — misses the signal",
      "Lesson: global features cannot replace local shape when timing matters"]),
    ("SHAP as a Design Tool", GREEN,
     ["Not just post-hoc — actively used to prune features before training",
      "Reduced feature space by 40–45% across all three tasks",
      "Faster inference, lower memory, more interpretable models",
      "Critical for deployment on edge devices and wearables"]),
]

for i, (title, color, bullets) in enumerate(panels):
    lx = 0.35 + i * 4.32
    add_rect(slide, lx, 1.35, 4.1, 5.55, WHITE)
    add_rect(slide, lx, 1.35, 4.1, 0.5, color)
    add_text(slide, title, lx + 0.1, 1.38, 3.9, 0.42,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bullets):
        add_text(slide, f"▸  {b}", lx + 0.12, 2.0 + j * 0.85, 3.85, 0.8,
                 font_size=12, color=TEXT_DARK, wrap=True)

# Challenges box
add_rect(slide, 0.35, 7.0, 12.6, 0.0, LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSION & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Conclusion & Future Work",
             "Key takeaways and what comes next")

# Left: conclusions
add_rect(slide, 0.35, 1.35, 6.3, 5.55, WHITE)
add_text(slide, "Key Conclusions", 0.5, 1.43, 6.0, 0.42,
         font_size=15, bold=True, color=MID_BLUE)

conclusions = [
    "Feature engineering still matters — Catch22 + LightGBM rivals deep learning at a fraction of the cost",
    "Lag + Rolling Statistics are the strongest features for time-series forecasting (MSE −82%)",
    "SHAP-guided pruning is a reliable, principled way to reduce dimensionality without hurting accuracy",
    "Self-supervised embeddings did not outperform explicit statistical features on ETTh1",
    "Cross-domain generalization is limited: global features fail when local timing is the discriminant",
]
for i, c in enumerate(conclusions):
    add_text(slide, f"▸  {c}", 0.5, 2.0 + i * 0.88, 6.0, 0.82,
             font_size=12.5, color=TEXT_DARK, wrap=True)

# Right: future work
add_rect(slide, 6.9, 1.35, 6.1, 5.55, WHITE)
add_text(slide, "Future Work", 7.05, 1.43, 5.8, 0.42,
         font_size=15, bold=True, color=ACCENT)

future = [
    ("LLM-Assisted Feature Suggestion",
     "Use LLMs to read dataset metadata and suggest domain-specific features automatically"),
    ("Edge Device Deployment",
     "Develop ultra-lightweight Catch22 variants that run on IoT sensors in real-time"),
    ("Real WESAD Data",
     "Validate pipeline on the actual WESAD wearable stress dataset (requires institutional access)"),
    ("XAI-First AutoML",
     "Make SHAP a primary optimization objective in AutoML, not an afterthought"),
]
for i, (title, desc) in enumerate(future):
    y = 2.0 + i * 1.2
    add_text(slide, f"◆  {title}", 7.05, y, 5.8, 0.4,
             font_size=13, bold=True, color=DARK_BLUE)
    add_text(slide, desc, 7.25, y + 0.4, 5.6, 0.65,
             font_size=11.5, color=TEXT_MID, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — REFERENCES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "References", "IEEE citation format")

refs = [
    "[1]  H. Badi, G. Forestier, and J. Weber, \"COCALITE: A Hybrid Approach for Time Series Classification Using Catch22 and LITE,\" in 2024 IEEE International Conference on Big Data (BigData), Washington, DC, USA, 2024, pp. 5234–5243.",
    "[2]  M. D'Aversa, F. Giannini, and M. Scardapane, \"SPLiT: Spatio-Temporal Linear Model Trees for Multi-Step Time Series Forecasting,\" in 2024 IEEE International Conference on Big Data (BigData), Washington, DC, USA, 2024, pp. 3412–3421.",
    "[3]  L. Valerio, A. Passarella, and M. Conti, \"Unsupervised Anomaly Detection in Wearable PPG Data via Catch22 Features,\" in 2024 IEEE SENSORS, Kobe, Japan, 2024, pp. 1–4.",
    "[4]  T. Okadome and J. Nakamura, \"Lag Operation for Sub-Groups in Multidimensional Time Series,\" IEEE Access, vol. 12, pp. 15432–15445, 2024.",
    "[5]  O. Petrosian, D. Wang, and S. Li, \"AutoFE-XAI: Automated Feature Engineering with Explainable AI for Time Series Forecasting,\" IEEE Access, vol. 13, pp. 2109–2120, 2025.",
    "[6]  C. H. Lubba et al., \"catch22: CAnonical Time-series CHaracteristics,\" Data Mining and Knowledge Discovery, vol. 33, no. 6, pp. 1821–1852, 2019.",
]

for i, ref in enumerate(refs):
    add_text(slide, ref, 0.4, 1.45 + i * 0.88, 12.5, 0.82,
             font_size=12, color=TEXT_DARK, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU / Q&A
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 0, 0.55, 7.5, ACCENT)
add_rect(slide, 0, 3.2, 13.33, 0.08, ACCENT)

add_text(slide, "Thank You", 0.85, 1.0, 11.5, 1.5,
         font_size=56, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "Questions & Discussion",
         0.85, 2.5, 11.5, 0.7,
         font_size=24, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)

add_text(slide, "GitHub Repository:",
         0.85, 3.55, 4.0, 0.4, font_size=13, color=RGBColor(0x88, 0xAA, 0xCC))
add_text(slide,
    "github.com/Ash-projects-personal/CSCE5222-Group18-FeatureEngineering",
    0.85, 3.95, 11.5, 0.45, font_size=14, bold=True, color=WHITE)

add_text(slide, "Ashish Rathnakar Shetty  ·  ashishrathnakarshetty@unt.edu",
         0.85, 4.65, 11.5, 0.4, font_size=13, color=RGBColor(0xB0, 0xC8, 0xE8))
add_text(slide, "Kushal Sai Venigalla  ·  KushalSaiVenigalla@unt.edu",
         0.85, 5.1, 11.5, 0.4, font_size=13, color=RGBColor(0xB0, 0xC8, 0xE8))

add_rect(slide, 0, 6.9, 13.33, 0.6, RGBColor(0x0D, 0x1F, 0x3C))
add_text(slide, "CSCE 5222 Feature Engineering  |  Group 18  |  University of North Texas",
         0.85, 6.95, 12, 0.45, font_size=12, color=RGBColor(0x88, 0xAA, 0xCC))


# ─── Save ────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
