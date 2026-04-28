"""
fix_slide16.py — fixes the Thank You slide (slide 16)
Clears the old template content placeholder and cleans up the layout.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT    = '/home/ubuntu/project2/slides/Group18_UNT_Presentation.pptx'
UNT_GREEN = RGBColor(0x00, 0x85, 0x3E)
BLACK     = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)

prs = Presentation(OUTPUT)
s = prs.slides[15]  # slide 16

# Step 1: Remove ALL non-placeholder shapes (the textboxes we added before)
to_remove = [shape for shape in s.shapes if not shape.is_placeholder]
for shape in to_remove:
    shape._element.getparent().remove(shape._element)

# Step 2: Clear the content placeholder (idx=1) which has old template text
for ph in s.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.clear()
        # Make it invisible by setting empty text
        tf.paragraphs[0].text = ''
        # Move it off-screen
        ph.left   = Inches(20)
        ph.top    = Inches(20)
        ph.width  = Inches(0.1)
        ph.height = Inches(0.1)

# Step 3: Add clean textboxes for the Thank You content
def add_tb(text, left, top, width, height, fs, bold=False,
           color=BLACK, align=PP_ALIGN.CENTER):
    txBox = s.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.color.rgb = color

add_tb("Questions & Discussion",
       1.5, 2.8, 10.33, 0.6, 26, False, DARK_GRAY)

# Divider line using a thin rectangle
shape = s.shapes.add_shape(
    1, Inches(3.5), Inches(3.55), Inches(6.33), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = UNT_GREEN
shape.line.fill.background()

add_tb("Ashish Rathnakar Shetty  ·  ashishrathnakarshetty@unt.edu",
       1.5, 3.75, 10.33, 0.48, 18, False, BLACK)

add_tb("Kushal Sai Venigalla  ·  KushalSaiVenigalla@unt.edu",
       1.5, 4.28, 10.33, 0.48, 18, False, BLACK)

add_tb("github.com/Ash-projects-personal/CSCE5222-Group18-FeatureEngineering",
       1.5, 4.95, 10.33, 0.45, 15, False, UNT_GREEN)

add_tb("CSCE 5222 Feature Engineering  |  Group 18  |  University of North Texas",
       1.5, 5.5, 10.33, 0.4, 14, False, DARK_GRAY)

prs.save(OUTPUT)
print("Slide 16 fixed and saved.")
