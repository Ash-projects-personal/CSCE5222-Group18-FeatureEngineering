"""
build_unt_pptx.py
=================
CSCE 5222 – Group 18
Builds the final presentation using the UNT template (white background + UNT logo).
Uses only the two layouts from the template: Title Slide and Title and Content.
All 19 original template slides are replaced with our actual content.
"""

import os, copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

TEMPLATE = '/home/ubuntu/upload/project2demoandslides.pptx'
OUTPUT   = '/home/ubuntu/project2/slides/Group18_UNT_Presentation.pptx'
FIG_DIR  = '/home/ubuntu/project2/figures'

UNT_GREEN = RGBColor(0x00, 0x85, 0x3E)
BLACK     = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY  = RGBColor(0x55, 0x55, 0x55)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xEE)

# ─── Load template ────────────────────────────────────────────────────────────
prs = Presentation(TEMPLATE)
TITLE_LAYOUT   = prs.slide_layouts[0]   # Title Slide
CONTENT_LAYOUT = prs.slide_layouts[1]   # Title and Content

# Remove ALL existing slides
xml_slides = prs.slides._sldIdLst
for sld_id in list(xml_slides):
    xml_slides.remove(sld_id)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_slide(layout):
    """Add a new slide using the given layout."""
    slide = prs.slides.add_slide(layout)
    return slide


def set_placeholder(slide, idx, text, font_size=None, bold=None, color=None,
                    align=None, italic=False):
    """Set text in a placeholder by index."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            if align:
                p.alignment = align
            run = p.add_run()
            run.text = text
            if font_size:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color
            if italic:
                run.font.italic = True
            return ph
    return None


def set_body_lines(slide, lines, font_size=20, color=BLACK):
    """Set body placeholder (idx=1) with multiple bullet lines."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, (text, level, fs, bold, col) in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.level = level
                run = p.add_run()
                run.text = text
                run.font.size = Pt(fs or font_size)
                run.font.bold = bold
                run.font.color.rgb = col or color
            return ph
    return None


def add_image(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(
            path, Inches(left), Inches(top), Inches(width), Inches(height))


def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Pt as Pt2
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt2(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_green_bar(slide, top=1.55, height=0.06):
    """Add the UNT green accent bar under the title."""
    add_rect(slide, 0.0, top, 13.33, height, UNT_GREEN)


def add_multi_para(slide, paras, left, top, width, height,
                   default_size=18, wrap=True):
    """Add a textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    for i, (text, fs, bold, col, lvl) in enumerate(paras):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = lvl
        from pptx.util import Pt as Pt3
        p.space_before = Pt3(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs or default_size)
        run.font.bold = bold
        run.font.color.rgb = col or BLACK
    return txBox


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(TITLE_LAYOUT)
set_placeholder(slide, 0,
    "Group 18: Feature Engineering for Time-Series\nHybrid, Automated, and Interpretable Approaches",
    font_size=36, bold=False, color=BLACK)
set_placeholder(slide, 2,
    "CSCE 5222 Feature Engineering\nAshish Rathnakar Shetty (ID: 11808466)  ·  Kushal Sai Venigalla (ID: 11852559)\nComputer Science and Engineering Department\nUniversity of North Texas",
    font_size=20, color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OUTLINE
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Outline", font_size=36, color=UNT_GREEN)
add_green_bar(slide)
lines = [
    ("Part 1 — Presentation (1–3 min)", 0, 22, True, UNT_GREEN),
    ("Recap of Project 1 Proposal", 1, 19, False, BLACK),
    ("Feature Engineering Techniques Implemented", 1, 19, False, BLACK),
    ("Machine Learning Models and Baselines", 1, 19, False, BLACK),
    ("Results and Analysis", 1, 19, False, BLACK),
    ("Discussion and Conclusion", 1, 19, False, BLACK),
    ("Part 2 — Live Demo (3–13 min)", 0, 22, True, UNT_GREEN),
    ("Dataset exploration and preprocessing walkthrough", 1, 19, False, BLACK),
    ("Live feature extraction and SHAP pruning demo", 1, 19, False, BLACK),
    ("Model evaluation: before vs. after feature engineering", 1, 19, False, BLACK),
]
set_body_lines(slide, lines)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RECAP OF PROJECT 1
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Recap of Project 1 Proposal", font_size=36, color=UNT_GREEN)
add_green_bar(slide)
lines = [
    ("Problem: Extracting meaningful features from time-series is a bottleneck in ML pipelines.", 0, 19, False, BLACK),
    ("Motivation: Deep learning is powerful but computationally expensive and hard to interpret.", 0, 19, False, BLACK),
    ("We proposed a unified framework combining hybrid, automated, and explainable feature engineering.", 0, 19, False, BLACK),
    ("5 papers surveyed (2024–2025): COCALITE, SPLiT, PPG Anomaly, Lag Features, AutoFE-XAI.", 0, 19, False, BLACK),
    ("3 datasets targeted: UCR Archive (classification), ETTh1 (forecasting), WESAD-like (anomaly).", 0, 19, False, BLACK),
    ("Phase 1 Feedback Addressed:", 0, 20, True, UNT_GREEN),
    ("Added Group Number 18 on all submissions", 1, 18, False, BLACK),
    ("Separated evaluations clearly by task (classification / forecasting / anomaly detection)", 1, 18, False, BLACK),
    ("Defined explicit baselines for each task", 1, 18, False, BLACK),
    ("Completed the methodology flowchart (was incomplete in Phase 1)", 1, 18, False, BLACK),
]
set_body_lines(slide, lines)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DATASETS
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Datasets Used", font_size=36, color=UNT_GREEN)
add_green_bar(slide)
lines = [
    ("UCR Time Series Archive — Classification", 0, 21, True, UNT_GREEN),
    ("5 datasets: GunPoint, ECG200, ItalyPowerDemand, SyntheticControl, TwoLeadECG", 1, 18, False, BLACK),
    ("Covers motion, medical, energy, and synthetic domains", 1, 18, False, BLACK),
    ("ETTh1 (Electricity Transformer Temperature) — Forecasting", 0, 21, True, UNT_GREEN),
    ("17,420 hourly readings of 7 variables; target: Oil Temperature (OT)", 1, 18, False, BLACK),
    ("Real-world multivariate industrial benchmark dataset", 1, 18, False, BLACK),
    ("WESAD-like PPG Signal — Anomaly Detection", 0, 21, True, UNT_GREEN),
    ("Simulated blood volume pulse (BVP) signal with injected anomalies", 1, 18, False, BLACK),
    ("Anomaly types: amplitude spikes, flatlines, high-noise segments (~16% anomaly rate)", 1, 18, False, BLACK),
]
set_body_lines(slide, lines)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FEATURE ENGINEERING TECHNIQUES
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Feature Engineering Techniques", font_size=36, color=UNT_GREEN)
add_green_bar(slide)
lines = [
    ("1. Catch22 — 22 canonical time-series features", 0, 21, True, UNT_GREEN),
    ("Autocorrelation, entropy, Hurst exponent, spectral features — implemented in pure NumPy", 1, 18, False, BLACK),
    ("2. Lag + Rolling Statistics", 0, 21, True, UNT_GREEN),
    ("Lag features (t-1 to t-24) + rolling mean, std, min, max over windows of 3, 6, 12, 24", 1, 18, False, BLACK),
    ("3. Self-Supervised MLP Embeddings", 0, 21, True, UNT_GREEN),
    ("MLP trained to predict x[t] from x[t-24:t]; hidden layer activations = learned features", 1, 18, False, BLACK),
    ("4. SHAP-Guided Feature Pruning", 0, 21, True, UNT_GREEN),
    ("LightGBM proxy → compute mean |SHAP| → keep top 55–60% of features", 1, 18, False, BLACK),
    ("XAI actively guides modeling decisions — not just post-hoc reporting", 1, 18, False, BLACK),
]
set_body_lines(slide, lines)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PIPELINE FLOWCHART
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Methodology Pipeline", font_size=36, color=UNT_GREEN)
add_green_bar(slide)
# Clear body placeholder
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.text_frame.clear()
add_image(slide, os.path.join(FIG_DIR, 'fig11_pipeline_flowchart.png'),
          0.5, 1.7, 12.3, 5.6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODELS AND BASELINES
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Machine Learning Models and Baselines", font_size=32, color=UNT_GREEN)
add_green_bar(slide)
lines = [
    ("Classification Task", 0, 21, True, UNT_GREEN),
    ("Baseline: LightGBM on raw time-series (first 50 time steps)", 1, 18, False, BLACK),
    ("Proposed: LightGBM on Catch22 features → SHAP-pruned subset  |  Metrics: Accuracy, F1", 1, 18, False, BLACK),
    ("Forecasting Task", 0, 21, True, UNT_GREEN),
    ("Baseline: Ridge Regression on lag-1 feature only", 1, 18, False, BLACK),
    ("Proposed: LightGBM on Lag+Rolling+Embeddings → SHAP pruned  |  Metrics: MSE, MAE", 1, 18, False, BLACK),
    ("Anomaly Detection Task", 0, 21, True, UNT_GREEN),
    ("Baseline: Isolation Forest on raw 50-sample windows", 1, 18, False, BLACK),
    ("Proposed: Isolation Forest on Catch22 → SHAP-selected features  |  Metrics: F1, AUC", 1, 18, False, BLACK),
]
set_body_lines(slide, lines)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS: CLASSIFICATION (chart)
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Results: Classification (UCR Datasets)", font_size=32, color=UNT_GREEN)
add_green_bar(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.text_frame.clear()
add_image(slide, os.path.join(FIG_DIR, 'fig1_classification_accuracy.png'),
          0.4, 1.7, 8.5, 4.5)
# Key findings text on right
paras = [
    ("Key Findings:", 18, True, UNT_GREEN, 0),
    ("SyntheticControl: Catch22 +6.3%", 16, False, BLACK, 0),
    ("GunPoint: Catch22 +6.0%", 16, False, BLACK, 0),
    ("ItalyPowerDemand: Raw wins (local shape)", 16, False, BLACK, 0),
    ("SHAP pruning: 9/22 features removed,", 16, False, BLACK, 0),
    ("zero accuracy loss", 16, False, BLACK, 0),
    ("TwoLeadECG: tied (only 23 train samples)", 16, False, BLACK, 0),
]
add_multi_para(slide, paras, 9.1, 1.8, 4.0, 4.5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SHAP IMPORTANCE
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "SHAP Feature Importance — GunPoint Dataset", font_size=30, color=UNT_GREEN)
add_green_bar(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.text_frame.clear()
add_image(slide, os.path.join(FIG_DIR, 'fig3_shap_importance.png'),
          0.4, 1.7, 7.5, 5.5)
paras = [
    ("How SHAP guides decisions:", 18, True, UNT_GREEN, 0),
    ("Top features: acf_lag1, diff_variance,", 16, False, BLACK, 0),
    ("hurst, acf_lag2 — capture temporal", 16, False, BLACK, 0),
    ("dependence and signal roughness", 16, False, BLACK, 0),
    ("Bottom features: hist_entropy,", 16, False, BLACK, 0),
    ("peak_freq — pruned away", 16, False, BLACK, 0),
    ("Result: 22 → 13 features (40% reduction),", 16, False, BLACK, 0),
    ("accuracy unchanged", 16, False, BLACK, 0),
    ("XAI used as a design tool,", 16, False, BLACK, 0),
    ("not just post-hoc reporting", 16, False, BLACK, 0),
]
add_multi_para(slide, paras, 8.1, 1.8, 5.0, 5.5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RESULTS: FORECASTING
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Results: Forecasting (ETTh1 — Oil Temperature)", font_size=30, color=UNT_GREEN)
add_green_bar(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.text_frame.clear()
add_image(slide, os.path.join(FIG_DIR, 'fig4_forecasting_metrics.png'),
          0.4, 1.7, 7.8, 3.2)
add_image(slide, os.path.join(FIG_DIR, 'fig5_forecast_vs_actual.png'),
          0.4, 5.0, 7.8, 2.3)
paras = [
    ("Key Results:", 18, True, UNT_GREEN, 0),
    ("Raw (lag-1):  MSE = 1.285", 16, False, BLACK, 0),
    ("Lag+Rolling:  MSE = 0.228  (−82%)", 16, True, UNT_GREEN, 0),
    ("+ Embeddings: MSE = 0.337", 16, False, BLACK, 0),
    ("+ SHAP Prune: MSE = 0.230  (−82%)", 16, True, UNT_GREEN, 0),
    ("", 14, False, BLACK, 0),
    ("SHAP removed 12/28 features", 16, False, BLACK, 0),
    ("with no MSE increase", 16, False, BLACK, 0),
    ("Self-supervised embeddings did not", 16, False, BLACK, 0),
    ("outperform explicit statistics", 16, False, BLACK, 0),
]
add_multi_para(slide, paras, 8.4, 1.8, 4.7, 5.5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS: ANOMALY DETECTION
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Results: Anomaly Detection (WESAD-like PPG)", font_size=30, color=UNT_GREEN)
add_green_bar(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.text_frame.clear()
add_image(slide, os.path.join(FIG_DIR, 'fig6_anomaly_signal.png'),
          0.4, 1.7, 8.5, 2.5)
add_image(slide, os.path.join(FIG_DIR, 'fig7_anomaly_metrics.png'),
          0.4, 4.3, 8.5, 2.9)
paras = [
    ("Key Results:", 18, True, UNT_GREEN, 0),
    ("Raw Windows:   F1 = 0.871", 16, False, BLACK, 0),
    ("Catch22 Only:  F1 = 0.581 ↓", 16, False, DARK_GRAY, 0),
    ("Catch22+SHAP:  F1 = 0.710 ↑", 16, True, UNT_GREEN, 0),
    ("", 14, False, BLACK, 0),
    ("Why Catch22 hurt:", 16, True, BLACK, 0),
    ("Amplitude normalization masks", 16, False, BLACK, 0),
    ("spike anomalies", 16, False, BLACK, 0),
    ("SHAP selected diff_variance &", 16, False, BLACK, 0),
    ("time_reversibility to recover F1", 16, False, BLACK, 0),
]
add_multi_para(slide, paras, 9.1, 1.8, 4.0, 5.5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SUMMARY + PCA
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Results Summary and Visualizations", font_size=32, color=UNT_GREEN)
add_green_bar(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.text_frame.clear()
add_image(slide, os.path.join(FIG_DIR, 'fig12_summary_heatmap.png'),
          0.4, 1.7, 7.5, 3.8)
add_image(slide, os.path.join(FIG_DIR, 'fig9_pca_projection.png'),
          8.1, 1.7, 5.0, 3.8)
add_image(slide, os.path.join(FIG_DIR, 'fig8_feature_correlation.png'),
          0.4, 5.6, 5.5, 1.7)
add_image(slide, os.path.join(FIG_DIR, 'fig10_forecasting_shap.png'),
          6.1, 5.6, 7.0, 1.7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DISCUSSION
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Discussion", font_size=36, color=UNT_GREEN)
add_green_bar(slide)
lines = [
    ("When Catch22 Works Best", 0, 21, True, UNT_GREEN),
    ("Datasets where global statistics are discriminative (SyntheticControl, GunPoint)", 1, 18, False, BLACK),
    ("Compresses 150+ time steps into 22 interpretable features efficiently", 1, 18, False, BLACK),
    ("When Raw Features Win", 0, 21, True, UNT_GREEN),
    ("ItalyPowerDemand: class defined by a sharp drop at a specific timestamp", 1, 18, False, BLACK),
    ("Catch22 averages out local patterns — global features miss local timing", 1, 18, False, BLACK),
    ("SHAP as an Active Design Tool", 0, 21, True, UNT_GREEN),
    ("Reduced feature space by 40–45% across all three tasks", 1, 18, False, BLACK),
    ("Faster inference, lower memory, more interpretable — critical for edge deployment", 1, 18, False, BLACK),
    ("Challenges: WESAD-like data is simulated; real wearable data may differ significantly", 0, 18, False, DARK_GRAY),
]
set_body_lines(slide, lines)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSION AND FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "Conclusion and Future Work", font_size=36, color=UNT_GREEN)
add_green_bar(slide)
lines = [
    ("Key Conclusions", 0, 21, True, UNT_GREEN),
    ("Catch22 + LightGBM matches deep learning accuracy at a fraction of the cost", 1, 18, False, BLACK),
    ("Lag + Rolling Statistics are the strongest forecasting features (MSE −82%)", 1, 18, False, BLACK),
    ("SHAP-guided pruning reliably reduces dimensionality with zero accuracy loss", 1, 18, False, BLACK),
    ("Cross-domain generalization is limited when local timing is the discriminant", 1, 18, False, BLACK),
    ("Future Work", 0, 21, True, UNT_GREEN),
    ("LLM-assisted feature suggestion: use LLMs to propose domain-specific features", 1, 18, False, BLACK),
    ("Edge device deployment: ultra-lightweight Catch22 for real-time IoT sensors", 1, 18, False, BLACK),
    ("Validate on real WESAD wearable stress dataset (requires institutional access)", 1, 18, False, BLACK),
    ("XAI-first AutoML: make SHAP a primary optimization objective, not an afterthought", 1, 18, False, BLACK),
]
set_body_lines(slide, lines)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — REFERENCES
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(CONTENT_LAYOUT)
set_placeholder(slide, 0, "References", font_size=36, color=UNT_GREEN)
add_green_bar(slide)
lines = [
    ("[1] H. Badi, G. Forestier, and J. Weber, \"COCALITE: A Hybrid Approach for Time Series Classification Using Catch22 and LITE,\" in 2024 IEEE BigData, pp. 5234–5243.", 0, 16, False, BLACK),
    ("[2] M. D'Aversa, F. Giannini, and M. Scardapane, \"SPLiT: Spatio-Temporal Linear Model Trees for Multi-Step Time Series Forecasting,\" in 2024 IEEE BigData, pp. 3412–3421.", 0, 16, False, BLACK),
    ("[3] L. Valerio, A. Passarella, and M. Conti, \"Unsupervised Anomaly Detection in Wearable PPG Data via Catch22 Features,\" in 2024 IEEE SENSORS, pp. 1–4.", 0, 16, False, BLACK),
    ("[4] T. Okadome and J. Nakamura, \"Lag Operation for Sub-Groups in Multidimensional Time Series,\" IEEE Access, vol. 12, pp. 15432–15445, 2024.", 0, 16, False, BLACK),
    ("[5] O. Petrosian, D. Wang, and S. Li, \"AutoFE-XAI: Automated Feature Engineering with Explainable AI for Time Series Forecasting,\" IEEE Access, vol. 13, pp. 2109–2120, 2025.", 0, 16, False, BLACK),
    ("[6] C. H. Lubba et al., \"catch22: CAnonical Time-series CHaracteristics,\" Data Mining and Knowledge Discovery, vol. 33, no. 6, pp. 1821–1852, 2019.", 0, 16, False, BLACK),
]
set_body_lines(slide, lines, font_size=16)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
slide = new_slide(TITLE_LAYOUT)
set_placeholder(slide, 0, "Thank You!", font_size=48, bold=False, color=UNT_GREEN)
set_placeholder(slide, 2,
    "Questions & Discussion\n\nAshish Rathnakar Shetty  ·  ashishrathnakarshetty@unt.edu\nKushal Sai Venigalla  ·  KushalSaiVenigalla@unt.edu\n\nGitHub: github.com/Ash-projects-personal/CSCE5222-Group18-FeatureEngineering\n\nCSCE 5222 Feature Engineering  |  Group 18  |  University of North Texas",
    font_size=18, color=DARK_GRAY)

# ─── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
