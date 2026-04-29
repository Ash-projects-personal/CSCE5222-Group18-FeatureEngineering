"""
build_12_slides.py
Clean 12-slide UNT presentation. No leftover template content. No mess.
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = '/home/ubuntu/upload/project2demoandslides.pptx'
OUTPUT   = '/home/ubuntu/project2/slides/Group18_UNT_Presentation.pptx'
FIG_DIR  = '/home/ubuntu/project2/figures'

shutil.copy2(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)

UNT_GREEN = RGBColor(0x00, 0x85, 0x3E)
BLACK     = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)


def clear_slide(slide):
    """Remove ALL non-placeholder shapes."""
    to_remove = [s for s in slide.shapes if not s.is_placeholder]
    for s in to_remove:
        s._element.getparent().remove(s._element)


def clear_all_text(slide):
    """Wipe text from all placeholders."""
    for ph in slide.placeholders:
        try:
            tf = ph.text_frame
            tf.clear()
            if tf.paragraphs:
                tf.paragraphs[0].text = ''
        except Exception:
            pass


def set_ph(slide, idx, lines):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, (text, fs, bold, color, level) in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.level = level
                p.space_before = Pt(3)
                run = p.add_run()
                run.text = text
                if fs:    run.font.size  = Pt(fs)
                if bold is not None: run.font.bold = bold
                if color: run.font.color.rgb = color
            return ph


def add_tb(slide, text, l, t, w, h, fs=16, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.color.rgb = color


def add_img(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


def green_line(slide):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(1.58), Inches(13.33), Inches(0.055))
    shape.fill.solid()
    shape.fill.fore_color.rgb = UNT_GREEN
    shape.line.fill.background()


def hide_body_ph(slide):
    """Move body placeholder off-screen so it doesn't show through."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            try:
                tf = ph.text_frame
                tf.clear()
                tf.paragraphs[0].text = ''
                ph.left = Inches(20); ph.top = Inches(20)
                ph.width = Inches(0.1); ph.height = Inches(0.1)
            except Exception:
                pass


sl = list(prs.slides)
# Template has 19 slides — we need only 12

# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
s = sl[0]
clear_slide(s)
clear_all_text(s)
set_ph(s, 0, [
    ("Group 18: Feature Engineering for Time-Series", 32, False, BLACK, 0),
    ("Hybrid, Automated, and Interpretable Approaches", 24, False, DARK_GRAY, 0),
])
add_tb(s, "CSCE 5222 Feature Engineering  |  Group 18",
       1.5, 3.9, 10.33, 0.5, 18, False, UNT_GREEN, PP_ALIGN.CENTER)
add_tb(s, "Ashish Rathnakar Shetty (ID: 11808466)  ·  Kushal Sai Venigalla (ID: 11852559)",
       1.5, 4.5, 10.33, 0.55, 20, True, BLACK, PP_ALIGN.CENTER)
add_tb(s, "Computer Science and Engineering Department  |  University of North Texas",
       1.5, 5.15, 10.33, 0.45, 16, False, DARK_GRAY, PP_ALIGN.CENTER)

# ── SLIDE 2: OUTLINE ──────────────────────────────────────────────────────────
s = sl[1]
clear_slide(s)
set_ph(s, 0, [("Outline", 36, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("Part 1 — Presentation", 22, True, UNT_GREEN, 0),
    ("Recap of Project 1 Proposal", 19, False, BLACK, 1),
    ("Datasets and Feature Engineering Techniques", 19, False, BLACK, 1),
    ("Machine Learning Models and Baselines", 19, False, BLACK, 1),
    ("Results: Classification, Forecasting, Anomaly Detection", 19, False, BLACK, 1),
    ("Discussion and Conclusion", 19, False, BLACK, 1),
    ("Part 2 — Live Demo (Google Colab Notebook)", 22, True, UNT_GREEN, 0),
    ("Walkthrough of the Catch22 implementation", 19, False, BLACK, 1),
    ("Live SHAP-guided pruning across all three tasks", 19, False, BLACK, 1),
])

# ── SLIDE 3: RECAP ────────────────────────────────────────────────────────────
s = sl[2]
clear_slide(s)
set_ph(s, 0, [("Recap of Project 1 Proposal", 34, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("Problem: Extracting meaningful features from time-series is a bottleneck.", 19, False, BLACK, 0),
    ("Motivation: Deep learning is powerful but slow and hard to interpret.", 19, False, BLACK, 0),
    ("We surveyed 5 papers (2024–2025): COCALITE, SPLiT, PPG Anomaly, Lag, AutoFE-XAI.", 19, False, BLACK, 0),
    ("Three datasets: UCR (classification), ETTh1 (forecasting), WESAD-like (anomaly).", 19, False, BLACK, 0),
    ("Phase 1 Feedback Addressed:", 21, True, UNT_GREEN, 0),
    ("Group Number 18 added to all submissions", 18, False, BLACK, 1),
    ("Evaluations clearly separated by task; explicit baselines defined", 18, False, BLACK, 1),
    ("Methodology flowchart completed (was incomplete in Phase 1)", 18, False, BLACK, 1),
])

# ── SLIDE 4: DATASETS + FEATURES ──────────────────────────────────────────────
s = sl[3]
clear_slide(s)
set_ph(s, 0, [("Datasets and Feature Engineering", 32, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("Datasets", 21, True, UNT_GREEN, 0),
    ("UCR (classification): GunPoint, ECG200, ItalyPowerDemand, SyntheticControl, TwoLeadECG", 17, False, BLACK, 1),
    ("ETTh1 (forecasting): 17,420 hourly readings, 7 sensor variables, target = Oil Temperature", 17, False, BLACK, 1),
    ("WESAD-like PPG (anomaly): 2,000-sample BVP signal with injected anomalies (15.9% rate)", 17, False, BLACK, 1),
    ("Four Feature Engineering Strategies", 21, True, UNT_GREEN, 0),
    ("Catch22 — 22 canonical features implemented from scratch in pure NumPy", 17, False, BLACK, 1),
    ("Lag + Rolling Statistics — lags 1 to 24 + rolling mean/std/min/max over 4 windows", 17, False, BLACK, 1),
    ("Self-Supervised MLP Embeddings — 32-unit MLP predicts x[t] from x[t-24:t]", 17, False, BLACK, 1),
    ("SHAP-Guided Pruning — keep top 55–60% of features by mean |SHAP| value", 17, False, BLACK, 1),
])

# ── SLIDE 5: PIPELINE ─────────────────────────────────────────────────────────
s = sl[4]
clear_slide(s)
set_ph(s, 0, [("Methodology Pipeline", 36, False, UNT_GREEN, 0)])
green_line(s)
hide_body_ph(s)
add_img(s, os.path.join(FIG_DIR, 'fig11_pipeline_flowchart.png'), 0.5, 1.7, 12.3, 5.6)

# ── SLIDE 6: MODELS & BASELINES ───────────────────────────────────────────────
s = sl[5]
clear_slide(s)
set_ph(s, 0, [("Machine Learning Models and Baselines", 30, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("Classification", 21, True, UNT_GREEN, 0),
    ("Baseline: GradientBoosting on raw time-series (first 50 steps)", 18, False, BLACK, 1),
    ("Proposed: GradientBoosting on Catch22 → SHAP-pruned  |  Metrics: Accuracy, F1", 18, False, BLACK, 1),
    ("Forecasting", 21, True, UNT_GREEN, 0),
    ("Baseline: Ridge Regression on lag-1 only", 18, False, BLACK, 1),
    ("Proposed: GradientBoosting on Lag+Rolling+Embeddings → SHAP pruned  |  MSE, MAE", 18, False, BLACK, 1),
    ("Anomaly Detection", 21, True, UNT_GREEN, 0),
    ("Baseline: Isolation Forest on raw 50-sample windows", 18, False, BLACK, 1),
    ("Proposed: Isolation Forest on Catch22 → SHAP-selected  |  F1, AUC", 18, False, BLACK, 1),
])

# ── SLIDE 7: CLASSIFICATION RESULTS ───────────────────────────────────────────
s = sl[6]
clear_slide(s)
set_ph(s, 0, [("Results: Classification (UCR Datasets)", 30, False, UNT_GREEN, 0)])
green_line(s)
hide_body_ph(s)
add_img(s, os.path.join(FIG_DIR, 'fig1_classification_accuracy.png'), 0.4, 1.7, 8.5, 4.5)
add_tb(s, "Key Findings:", 9.1, 1.8, 4.0, 0.4, 17, True, UNT_GREEN)
findings = [
    ("SyntheticControl: SHAP +9.3%", True, UNT_GREEN),
    ("TwoLeadECG: Catch22 +8.8%", True, UNT_GREEN),
    ("ECG200: Catch22 +5.0%", True, UNT_GREEN),
    ("ItalyPowerDemand: Raw wins", False, DARK_GRAY),
    ("(local timing event missed", False, DARK_GRAY),
    ("by global statistics)", False, DARK_GRAY),
    ("SHAP pruning: 9/22 removed,", False, BLACK),
    ("zero accuracy loss", False, BLACK),
]
for i, (t, b, c) in enumerate(findings):
    add_tb(s, f"• {t}", 9.1, 2.3 + i * 0.45, 4.0, 0.42, 14, b, c)

# ── SLIDE 8: FORECASTING RESULTS ──────────────────────────────────────────────
s = sl[7]
clear_slide(s)
set_ph(s, 0, [("Results: Forecasting (ETTh1 — Oil Temperature)", 28, False, UNT_GREEN, 0)])
green_line(s)
hide_body_ph(s)
add_img(s, os.path.join(FIG_DIR, 'fig4_forecasting_metrics.png'), 0.4, 1.7, 7.8, 3.0)
add_img(s, os.path.join(FIG_DIR, 'fig5_forecast_vs_actual.png'), 0.4, 4.8, 7.8, 2.5)
add_tb(s, "Key Results:", 8.4, 1.8, 4.7, 0.4, 17, True, UNT_GREEN)
results = [
    ("Raw (lag-1):     MSE = 1.2856", False, BLACK),
    ("Lag+Rolling:     MSE = 0.4836", True, UNT_GREEN),
    ("                 (−62%)", False, UNT_GREEN),
    ("+ Embeddings:    MSE = 0.4807", False, BLACK),
    ("+ SHAP Prune:    MSE = 0.4734", True, UNT_GREEN),
    ("                 (best)", False, UNT_GREEN),
    ("", False, BLACK),
    ("SHAP removed 12/28 features", False, BLACK),
    ("Embeddings did not help", False, DARK_GRAY),
]
for i, (t, b, c) in enumerate(results):
    add_tb(s, t, 8.4, 2.3 + i * 0.45, 4.7, 0.42, 14, b, c)

# ── SLIDE 9: ANOMALY DETECTION RESULTS ────────────────────────────────────────
s = sl[8]
clear_slide(s)
set_ph(s, 0, [("Results: Anomaly Detection (WESAD-like PPG)", 28, False, UNT_GREEN, 0)])
green_line(s)
hide_body_ph(s)
add_img(s, os.path.join(FIG_DIR, 'fig6_anomaly_signal.png'), 0.4, 1.7, 8.5, 2.5)
add_img(s, os.path.join(FIG_DIR, 'fig7_anomaly_metrics.png'), 0.4, 4.3, 8.5, 2.9)
add_tb(s, "Key Results:", 9.1, 1.8, 4.0, 0.4, 17, True, UNT_GREEN)
anom = [
    ("Raw Windows:   F1 = 0.871", False, BLACK),
    ("Catch22 Only:  F1 = 0.613 ↓", False, DARK_GRAY),
    ("Catch22+SHAP:  F1 = 0.645 ↑", True, UNT_GREEN),
    ("", False, BLACK),
    ("Why Catch22 hurt:", True, BLACK),
    ("Normalization removes", False, BLACK),
    ("amplitude information", False, BLACK),
    ("Why SHAP helped:", True, BLACK),
    ("Kept diff_variance —", False, BLACK),
    ("captures roughness", False, BLACK),
]
for i, (t, b, c) in enumerate(anom):
    add_tb(s, t, 9.1, 2.3 + i * 0.42, 4.0, 0.4, 13, b, c)

# ── SLIDE 10: DISCUSSION ──────────────────────────────────────────────────────
s = sl[9]
clear_slide(s)
set_ph(s, 0, [("Discussion", 36, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("When Catch22 Works Best", 21, True, UNT_GREEN, 0),
    ("Datasets where global statistics differ between classes (SyntheticControl, ECG200, TwoLeadECG)", 18, False, BLACK, 1),
    ("Compresses 150+ time steps into 22 interpretable features; acts as regularization", 18, False, BLACK, 1),
    ("When Raw Features Win", 21, True, UNT_GREEN, 0),
    ("ItalyPowerDemand: class boundary is a sharp drop at one specific hour", 18, False, BLACK, 1),
    ("Catch22 averages over the whole signal, missing local timing events", 18, False, BLACK, 1),
    ("SHAP as a Design Tool", 21, True, UNT_GREEN, 0),
    ("Removed 40–45% of features across all three tasks with zero accuracy loss", 18, False, BLACK, 1),
    ("Faster inference, smaller model — important for edge deployment", 18, False, BLACK, 1),
    ("Limitation: WESAD-like data is simulated; real wearable data may differ", 18, False, DARK_GRAY, 0),
])

# ── SLIDE 11: CONCLUSION ──────────────────────────────────────────────────────
s = sl[10]
clear_slide(s)
set_ph(s, 0, [("Conclusion and Future Work", 32, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("Key Conclusions", 21, True, UNT_GREEN, 0),
    ("Catch22 + GradientBoosting beats raw features on 3/5 UCR datasets (best: +9.3% SyntheticControl)", 17, False, BLACK, 1),
    ("Lag + Rolling Statistics reduce forecasting MSE by 62%; SHAP pruning achieves best MSE of 0.4734", 17, False, BLACK, 1),
    ("SHAP pruning removes 40–45% of features across all tasks with zero accuracy loss", 17, False, BLACK, 1),
    ("Catch22 fails when class differences are local timing events, not global patterns", 17, False, BLACK, 1),
    ("Future Work", 21, True, UNT_GREEN, 0),
    ("Build a lighter Catch22 variant with 10–15 features focused on what SHAP keeps most", 17, False, BLACK, 1),
    ("Validate the anomaly detection pipeline on the real WESAD wearable dataset", 17, False, BLACK, 1),
    ("Explore LLM-assisted feature suggestion based on dataset metadata", 17, False, BLACK, 1),
])

# ── SLIDE 12: REFERENCES ──────────────────────────────────────────────────────
s = sl[11]
clear_slide(s)
set_ph(s, 0, [("References", 36, False, UNT_GREEN, 0)])
green_line(s)
set_ph(s, 1, [
    ("[1] H. Badi, G. Forestier, J. Weber, \"COCALITE: A Hybrid Approach for Time Series Classification Using Catch22 and LITE,\" IEEE BigData, 2024, pp. 5234–5243.", 13, False, BLACK, 0),
    ("[2] M. D'Aversa, F. Giannini, M. Scardapane, \"SPLiT: Spatio-Temporal Linear Model Trees for Multi-Step Time Series Forecasting,\" IEEE BigData, 2024, pp. 3412–3421.", 13, False, BLACK, 0),
    ("[3] L. Valerio, A. Passarella, M. Conti, \"Unsupervised Anomaly Detection in Wearable PPG Data via Catch22 Features,\" IEEE SENSORS, 2024, pp. 1–4.", 13, False, BLACK, 0),
    ("[4] T. Okadome, J. Nakamura, \"Lag Operation for Sub-Groups in Multidimensional Time Series,\" IEEE Access, vol. 12, pp. 15432–15445, 2024.", 13, False, BLACK, 0),
    ("[5] O. Petrosian, D. Wang, S. Li, \"AutoFE-XAI: Automated Feature Engineering with Explainable AI for Time Series Forecasting,\" IEEE Access, vol. 13, pp. 2109–2120, 2025.", 13, False, BLACK, 0),
    ("[6] C. H. Lubba et al., \"catch22: CAnonical Time-series CHaracteristics,\" Data Mining and Knowledge Discovery, vol. 33, no. 6, pp. 1821–1852, 2019.", 13, False, BLACK, 0),
    ("[7] H. A. Dau et al., \"The UCR Time Series Archive,\" IEEE/CAA Journal of Automatica Sinica, vol. 6, no. 6, pp. 1293–1305, 2019.", 13, False, BLACK, 0),
    ("[8] H. Zhou et al., \"Informer: Beyond Efficient Transformer for Long Sequence Time-Series Forecasting,\" AAAI, 2021, pp. 11106–11115.", 13, False, BLACK, 0),
    ("[9] P. Schmidt et al., \"Introducing WESAD, a Multimodal Dataset for Wearable Stress and Affect Detection,\" ACM ICMI, 2018, pp. 400–408.", 13, False, BLACK, 0),
])

# ── REMOVE SLIDES 13-19 (template extras) ─────────────────────────────────────
xml_slides = prs.slides._sldIdLst
all_ids = list(xml_slides)
for sld_id in all_ids[12:]:
    xml_slides.remove(sld_id)

print(f"Final slide count: {len(prs.slides)}")
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
